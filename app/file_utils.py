"""Utility functions to list files and extract text from PDF, DOCX, PPTX"""

from pathlib import Path
from PyPDF2 import PdfReader

## added packages (could use OCR for better text-extraction)
from docx import Document
from pptx import Presentation

# renamed "files"
FILE_FOLDER = Path(__file__).parent.parent / "files"
FILE_FOLDER.mkdir(exist_ok=True)


def list_files():
    """Return list of supported filenames (PDF, DOCX, PPTX)"""
    exts = ["*.pdf", "*.docx", "*.pptx"]
    files = []
    for ext in exts:
        files.extend(FILE_FOLDER.glob(ext))
    return sorted([f.name for f in files])


def extract_text(filename, max_chars=None):
    """
    Extract text from a file
    
    Args:
        filename (str): Name of the file in the `files` folder
        max_chars (int, optional): Maximum number of characters to return

    Returns:
        str: Extracted text or an error message
    """
    file_path = FILE_FOLDER / filename
    try:
        if filename.endswith(".pdf"):
            reader = PdfReader(file_path)
            # chunking by page
            text = ""
            for page in reader.pages:
                page_text = page.extract_text() or ""
                if page_text:
                    text += page_text + "\n"

        elif filename.endswith(".docx"):
            doc = Document(file_path)
            # chunking by paragraph
            text = "\n".join(p.text for p in doc.paragraphs)

        elif filename.endswith(".pptx"):
            prs = Presentation(file_path)
            # chunking by slide
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    # use getattr() to avoid BaseShape attribute errors
                    txt = getattr(shape, "text", "")
                    if txt and txt.strip():
                        text += txt + "\n"

        else:
            return "Unsupported file type"

        text = text.strip() or f"No text found in {filename}"
        # truncate if max_chars provided
        if max_chars:
            text = text[:max_chars]
        return text

    except Exception as e:
        return f"Error reading file: {e}"
